import json
import shutil
import subprocess
from pathlib import Path
from typing import Any

from jinja2 import Environment, FileSystemLoader
from pptx import Presentation
from pptx.util import Inches, Pt

from app.core.config import settings


class DocGenerator:
    def __init__(self):
        self.template_dir = Path(settings.template_dir)
        self.template_dir.mkdir(parents=True, exist_ok=True)
        self.jinja = Environment(loader=FileSystemLoader(str(self.template_dir)))

    async def generate_html(self, project_data: dict[str, Any], level: str = "L1") -> str:
        template_name = f"{level.lower()}_template.html"
        try:
            template = self.jinja.get_template(template_name)
        except Exception:
            template = self.jinja.from_string(self._default_html_template(level))

        html = template.render(**project_data)
        output_path = Path(settings.upload_dir) / project_data.get("project_id", "unknown") / f"{level}_preview.html"
        output_path.parent.mkdir(parents=True, exist_ok=True)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(html)
        return str(output_path)

    async def generate_ppt(self, project_data: dict[str, Any], level: str = "L1") -> str:
        template_path = Path(settings.ppt_template_path)
        if template_path.exists():
            prs = Presentation(str(template_path))
        else:
            prs = self._create_default_ppt_template()

        # 简化映射：在第一页添加标题
        if prs.slides:
            slide = prs.slides[0]
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text:
                    shape.text_frame.text = project_data.get("title", "美陈设计方案")

        output_path = Path(settings.upload_dir) / project_data.get("project_id", "unknown") / f"{level}_方案.pptx"
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        return str(output_path)

    async def generate_pdf(self, project_data: dict[str, Any], level: str = "L1") -> str:
        html_path = await self.generate_html(project_data, level)
        output_path = Path(settings.upload_dir) / project_data.get("project_id", "unknown") / f"{level}_方案.pdf"
        output_path.parent.mkdir(parents=True, exist_ok=True)

        # 尝试 WeasyPrint
        try:
            from weasyprint import HTML
            HTML(filename=html_path).write_pdf(str(output_path))
            return str(output_path)
        except Exception:
            pass

        # 降级：LibreOffice 从 PPT 转换
        try:
            ppt_path = await self.generate_ppt(project_data, level)
            subprocess.run(
                [
                    "soffice",
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    str(output_path.parent),
                    str(ppt_path),
                ],
                check=True,
                timeout=60,
            )
            return str(output_path)
        except Exception:
            pass

        # 最终降级：HTML 直接重命名
        shutil.copy(html_path, output_path.with_suffix(".html"))
        return str(output_path.with_suffix(".html"))

    def _default_html_template(self, level: str) -> str:
        return """<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>{{ title }}</title>
    <style>
        body { font-family: "PingFang SC", "Microsoft YaHei", sans-serif; margin: 40px; color: #333; }
        h1 { color: #2E5C8A; border-bottom: 3px solid #2E5C8A; padding-bottom: 10px; }
        h2 { color: #4A4A4A; margin-top: 30px; }
        .section { margin-bottom: 30px; }
        .image-grid { display: grid; grid-template-columns: repeat(3, 1fr); gap: 15px; }
        .image-grid img { width: 100%; border-radius: 8px; }
        table { width: 100%; border-collapse: collapse; margin-top: 15px; }
        th, td { border: 1px solid #ddd; padding: 10px; text-align: left; }
        th { background: #f5f5f5; }
        @media print { body { margin: 20px; } }
    </style>
</head>
<body>
    <h1>{{ title }}</h1>
    <div class="section">
        <h2>设计故事</h2>
        <p>{{ story.story if story else "" }}</p>
    </div>
    <div class="section">
        <h2>氛围描述</h2>
        <p>{{ atmosphere.paragraph if atmosphere else "" }}</p>
    </div>
    <div class="section">
        <h2>参考图</h2>
        <div class="image-grid">
            {% for img in images %}<img src="{{ img }}" alt="参考图">{% endfor %}
        </div>
    </div>
    <div class="section">
        <h2>物料清单</h2>
        <table>
            <tr><th>名称</th><th>材质</th><th>规格</th><th>数量</th><th>单价</th><th>总价</th></tr>
            {% for item in material_list %}<tr>
                <td>{{ item.name }}</td><td>{{ item.material }}</td><td>{{ item.size }}</td>
                <td>{{ item.quantity }}</td><td>{{ item.unit_price }}</td><td>{{ item.total_price }}</td>
            </tr>{% endfor %}
        </table>
    </div>
    <div class="section">
        <h2>预算汇总</h2>
        <p>材料费：{{ budget.material_cost }} 元</p>
        <p>制作费：{{ budget.production_cost }} 元</p>
        <p>安装费：{{ budget.installation_cost }} 元</p>
        <p>设计费：{{ budget.design_fee }} 元</p>
        <p><strong>总计：{{ budget.total }} 元</strong></p>
    </div>
</body>
</html>"""

    def _create_default_ppt_template(self) -> Presentation:
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(blank_layout)

        # 添加标题文本框
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = title_box.text_frame
        tf.text = "美陈设计方案"
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        return prs


doc_generator = DocGenerator()
