from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from pptx import Presentation

from app.services.llm_client import llm_client


class PPTParser:
    async def parse(self, file_path: Path | str) -> dict[str, Any]:
        path = Path(file_path)
        prs = Presentation(str(path))

        slides_text = []
        slide_images = []
        img_dir = path.parent / "ppt_images"
        img_dir.mkdir(exist_ok=True)

        for slide_idx, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            slide_text = "\n".join(texts)
            slides_text.append({"index": slide_idx + 1, "text": slide_text})

            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    image = shape.image
                    ext = image.ext
                    img_path = img_dir / f"slide{slide_idx + 1}_img.{ext}"
                    with open(img_path, "wb") as f:
                        f.write(image.blob)
                    slide_images.append(str(img_path))

        all_text = "\n".join(s["text"] for s in slides_text)
        structured = await self._extract_requirements(all_text)

        return {
            "source_type": "ppt",
            "file_path": str(path),
            "slide_count": len(prs.slides),
            "slides": slides_text,
            "images": slide_images,
            **structured,
        }

    async def _extract_requirements(self, text: str) -> dict[str, Any]:
        system = "你是一位需求文档分析师，擅长从 PPT 方案中提取结构化信息。"
        prompt = (
            f"请从以下 PPT 内容中提取关键信息，输出 JSON（只输出 JSON）：\n\n{text[:4000]}\n\n"
            "输出字段：\n"
            '{"theme": "主题", "style": "风格", "space_type": "空间类型", '
            '"budget": "预算", "timeline": "工期", "requirements": ["需求列表"], '
            '"restrictions": ["限制条件"]}'
        )
        raw = await llm_client.complete(system, prompt, json_mode=True)
        try:
            return json.loads(raw)
        except json.JSONDecodeError:
            return {"raw_description": raw, "parse_error": True}
